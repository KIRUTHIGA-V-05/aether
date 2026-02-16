from pptx import Presentation

def generate_ppt_file(history):
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Aether Lecture Session"
    title_slide.placeholders[1].text = "Automated Board Summary"

    for entry in history:
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = entry["intent"].replace("_", " ").title()
        slide.placeholders[1].text = entry["text"]

    output = "aether_summary.pptx"
    prs.save(output)
    return output
